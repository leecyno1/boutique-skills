#!/usr/bin/env python3
"""Cross-platform file-backed PPTX bridge for Scientific Illustrator.

The preferred Windows Microsoft PowerPoint backend edits the live COM model.
This bridge covers Microsoft PowerPoint for Mac and WPS Presentation on Windows
or macOS. It builds native editable OOXML objects with python-pptx, keeps them
in an isolated working copy, and reopens that copy in the selected presentation
application at explicit or sequence checkpoints. Background refresh preserves
the user's current foreground application by default. The system mouse and
keyboard are never used.
"""

from __future__ import annotations

import base64
import copy
import csv
import io
import json
import math
import os
from pathlib import Path
import plistlib
import re
import shutil
import subprocess
import sys
import tempfile
import time
import uuid

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


STATE_DIR = Path(
    os.environ.get("SCIENTIFIC_ILLUSTRATOR_STATE_DIR")
    or os.environ.get("SCIENTIFIC_ILLUSTRATOR_MAC_DIR")  # Backward compatibility with <= 1.5.2.
    or str(Path.home() / ".codex" / "scientific-illustrator" / "presentations")
).expanduser()
STATE_PATH = STATE_DIR / "session.json"
MAC_HOSTS = {
    "powerpoint": {
        "bundle_id": "com.microsoft.Powerpoint",
        "executable_name": "Microsoft PowerPoint",
        "paths": [Path("/Applications/Microsoft PowerPoint.app")],
    },
    "wps": {
        "bundle_id": "com.kingsoft.wpsoffice.mac",
        "executable_name": "wpsoffice",
        "paths": [
            Path("/Applications/wpsoffice.localized/wpsoffice.app"),
            Path("/Applications/wpsoffice.app"),
            Path("/Applications/WPS Office.app"),
            Path("/Applications/WPS.app"),
        ],
    },
}
MAC_HOST_ENV = {
    "powerpoint": "POWERPOINT_PRESENTATION_PATH",
    "wps": "WPS_PRESENTATION_PATH",
}


def _requested_host(args: dict | None = None, state: dict | None = None) -> str:
    explicit = (args or {}).get("host_application")
    configured = os.environ.get("SCIENTIFIC_ILLUSTRATOR_PPT_HOST", "auto")
    saved = (state or {}).get("host_application")
    value = str(explicit or saved or configured or "auto").strip().lower()
    if value not in {"auto", "powerpoint", "wps"}:
        raise ValueError("host_application must be auto, powerpoint, or wps")
    return value


def _windows_wps_candidates() -> list[Path]:
    candidates: list[Path] = []
    configured = os.environ.get("WPS_PRESENTATION_PATH")
    if configured:
        configured_path = Path(configured).expanduser()
        candidates.append(configured_path)
        if configured_path.is_dir():
            for executable_name in ("wpp.exe", "wpsoffice.exe"):
                candidates.append(configured_path / executable_name)
                candidates.append(configured_path / "office6" / executable_name)
                candidates.extend(configured_path.glob(f"*/office6/{executable_name}"))

    for executable_name in ("wpp.exe", "wpsoffice.exe"):
        discovered = shutil.which(executable_name)
        if discovered:
            candidates.append(Path(discovered))

    try:
        import winreg

        for executable_name in ("wpp.exe", "wpsoffice.exe"):
            key_name = rf"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\{executable_name}"
            for hive in (winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE):
                access_modes = {winreg.KEY_READ}
                for flag_name in ("KEY_WOW64_32KEY", "KEY_WOW64_64KEY"):
                    flag = getattr(winreg, flag_name, 0)
                    if flag:
                        access_modes.add(winreg.KEY_READ | flag)
                for access in access_modes:
                    try:
                        with winreg.OpenKey(hive, key_name, 0, access) as key:
                            value, _ = winreg.QueryValueEx(key, None)
                        if value:
                            candidates.append(Path(str(value).strip('"')))
                    except OSError:
                        continue
    except ImportError:
        pass

    for variable in ("ProgramFiles", "ProgramFiles(x86)", "LOCALAPPDATA"):
        root = os.environ.get(variable)
        if not root:
            continue
        base = Path(root)
        for product_root in (base / "Kingsoft").glob("WPS Office*"):
            for executable_name in ("wpp.exe", "wpsoffice.exe"):
                candidates.append(product_root / "office6" / executable_name)
                candidates.extend(product_root.glob(f"*/office6/{executable_name}"))

    unique: list[Path] = []
    seen: set[str] = set()
    for candidate in candidates:
        key = os.path.normcase(os.path.abspath(str(candidate)))
        if key not in seen:
            seen.add(key)
            unique.append(candidate)
    return unique


def _mac_app_bundle(path: Path) -> Path | None:
    """Normalize either an app bundle or its main executable to the .app path."""
    candidate = path.expanduser()
    for item in (candidate, *candidate.parents):
        if item.name.lower().endswith(".app"):
            return item
    return None


def _mac_bundle_metadata(app_path: Path) -> dict:
    plist_path = app_path / "Contents" / "Info.plist"
    try:
        with plist_path.open("rb") as handle:
            data = plistlib.load(handle)
    except (OSError, plistlib.InvalidFileException):
        data = {}
    return {
        "bundle_id": data.get("CFBundleIdentifier"),
        "version": data.get("CFBundleShortVersionString") or data.get("CFBundleVersion") or "",
        "executable_name": data.get("CFBundleExecutable"),
    }


def _mac_host_candidates(host_name: str) -> list[Path]:
    definition = MAC_HOSTS[host_name]
    raw_candidates: list[Path] = []
    configured = os.environ.get(MAC_HOST_ENV[host_name])
    if configured:
        raw_candidates.append(Path(configured))
    raw_candidates.extend(definition["paths"])
    if shutil.which("mdfind"):
        query = f"kMDItemCFBundleIdentifier == '{definition['bundle_id']}'"
        discovered = _run(["mdfind", query], check=False, timeout=15)
        raw_candidates.extend(Path(line) for line in discovered.stdout.splitlines() if line.strip())

    result: list[Path] = []
    seen: set[str] = set()
    for candidate in raw_candidates:
        app_path = _mac_app_bundle(candidate)
        if app_path is None or not app_path.exists():
            continue
        key = os.path.realpath(str(app_path))
        if key in seen:
            continue
        metadata = _mac_bundle_metadata(app_path)
        if metadata["bundle_id"] and metadata["bundle_id"] != definition["bundle_id"]:
            continue
        seen.add(key)
        result.append(app_path)
    return result


def _mac_host_descriptor(host_name: str) -> dict:
    definition = MAC_HOSTS[host_name]
    for found in _mac_host_candidates(host_name):
        metadata = _mac_bundle_metadata(found)
        executable_name = metadata["executable_name"] or definition["executable_name"]
        executable_path = found / "Contents" / "MacOS" / executable_name
        if executable_path.is_file():
            return {
                "installed": True,
                "path": str(found),
                "bundle_id": metadata["bundle_id"] or definition["bundle_id"],
                "executable_path": str(executable_path),
                "version": metadata["version"],
            }
    return {
        "installed": False,
        "path": None,
        "bundle_id": definition["bundle_id"],
        "executable_path": None,
        "version": "",
    }


def _available_hosts() -> dict[str, dict]:
    if sys.platform == "darwin":
        return {key: _mac_host_descriptor(key) for key in MAC_HOSTS}
    if sys.platform == "win32":
        wps = next((item for item in _windows_wps_candidates() if item.is_file()), None)
        powerpoint_installed = False
        powershell = shutil.which("powershell.exe") or shutil.which("powershell")
        if powershell:
            probe = _run([powershell, "-NoProfile", "-Command", "if (Test-Path 'Registry::HKEY_CLASSES_ROOT\\PowerPoint.Application\\CLSID') { 'yes' }"], check=False)
            powerpoint_installed = "yes" in probe.stdout.lower()
        return {
            "powerpoint": {"installed": powerpoint_installed, "path": None, "bundle_id": None, "executable_path": None, "version": ""},
            "wps": {"installed": wps is not None, "path": str(wps) if wps else None, "bundle_id": None, "executable_path": str(wps) if wps else None, "version": ""},
        }
    empty = {"installed": False, "path": None, "bundle_id": None, "executable_path": None, "version": ""}
    return {"powerpoint": dict(empty), "wps": dict(empty)}


def _select_host(args: dict | None = None, state: dict | None = None, available: dict | None = None) -> tuple[str, dict]:
    requested = _requested_host(args, state)
    available = available or _available_hosts()
    if requested == "auto":
        for candidate in ("powerpoint", "wps"):
            if available[candidate]["installed"]:
                return candidate, available[candidate]
        return "powerpoint", available["powerpoint"]
    return requested, available[requested]


def _state() -> dict:
    if STATE_PATH.exists():
        try:
            return json.loads(STATE_PATH.read_text("utf-8"))
        except Exception:
            pass
    return {"path": None, "source_path": None, "metadata": {}, "owned": False, "refresh_pending": False, "last_refresh": None}


def _write_state(state: dict) -> None:
    STATE_DIR.mkdir(parents=True, exist_ok=True)
    STATE_PATH.write_text(json.dumps(state, ensure_ascii=False, indent=2), "utf-8")


def _run(command: list[str], *, check: bool = True, timeout: int = 120) -> subprocess.CompletedProcess:
    return subprocess.run(command, check=check, capture_output=True, text=True, timeout=timeout)


def _osascript(source: str, *arguments: str, check: bool = False) -> str:
    result = _run(["osascript", "-e", source, *arguments], check=check, timeout=30)
    return result.stdout.strip()


def _focus_policy() -> str:
    value = os.environ.get("SCIENTIFIC_ILLUSTRATOR_FOCUS_POLICY", "preserve").strip().lower()
    return value if value in {"preserve", "foreground"} else "preserve"


def _open_windows_presentation(file_path: Path, executable: str | None, focus_policy: str) -> dict:
    if focus_policy == "foreground":
        if executable:
            subprocess.Popen([executable, str(file_path)], close_fds=True)
            method = "direct-executable"
        else:
            os.startfile(str(file_path))
            method = "os-startfile"
        return {"open_dispatched": True, "dispatch_method": method, "dispatch_return_code": None}

    # SW_SHOWNOACTIVATE prevents a newly created WPS/PowerPoint window from
    # taking focus. Some already-running hosts ignore that hint, so restore the
    # previously focused window after dispatching the file-open request.
    import ctypes

    user32 = ctypes.windll.user32
    previous_foreground = user32.GetForegroundWindow()
    if executable:
        startupinfo = subprocess.STARTUPINFO()
        startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
        startupinfo.wShowWindow = 4  # SW_SHOWNOACTIVATE
        subprocess.Popen([executable, str(file_path)], close_fds=True, startupinfo=startupinfo)
        method = "direct-executable-no-activate"
    else:
        result = ctypes.windll.shell32.ShellExecuteW(None, "open", str(file_path), None, None, 4)
        if result <= 32:
            raise RuntimeError(f"Unable to open presentation in the background; ShellExecuteW returned {result}.")
        method = "shell-execute-no-activate"
    if previous_foreground:
        time.sleep(0.15)
        user32.SetForegroundWindow(previous_foreground)
    return {"open_dispatched": True, "dispatch_method": method, "dispatch_return_code": None}


def _parse_macos_process_ids(output: str, executable_path: str | None) -> list[int]:
    if not executable_path:
        return []
    expected = os.path.realpath(executable_path)
    process_ids = []
    for line in output.splitlines():
        fields = line.strip().split(maxsplit=1)
        if len(fields) != 2 or not fields[0].isdigit():
            continue
        if os.path.realpath(fields[1].strip()) == expected:
            process_ids.append(int(fields[0]))
    return process_ids


def _parse_windows_tasklist_process_ids(output: str, expected_image_name: str) -> list[int]:
    process_ids: list[int] = []
    for fields in csv.reader(io.StringIO(output)):
        if len(fields) < 2:
            continue
        image_name, process_id = fields[0].strip(), fields[1].strip()
        if image_name.lower() == expected_image_name.lower() and process_id.isdigit():
            process_ids.append(int(process_id))
    return process_ids


def _main_process_ids(host_name: str, host: dict) -> list[int]:
    if not host.get("installed"):
        return []
    if sys.platform == "darwin":
        process_list = _run(["ps", "-ww", "-axo", "pid=,comm="], check=False, timeout=15)
        return _parse_macos_process_ids(process_list.stdout, host.get("executable_path"))
    if sys.platform == "win32":
        image_names = ["POWERPNT.EXE"] if host_name == "powerpoint" else ["wpp.exe", "wpsoffice.exe"]
        process_ids: list[int] = []
        for image_name in image_names:
            tasklist = _run(["tasklist", "/FI", f"IMAGENAME eq {image_name}", "/FO", "CSV", "/NH"], check=False)
            process_ids.extend(_parse_windows_tasklist_process_ids(tasklist.stdout, image_name))
        return sorted(set(process_ids))
    return []


def _parse_lsof_process_ids(output: str) -> set[int]:
    return {int(line[1:]) for line in output.splitlines() if line.startswith("p") and line[1:].isdigit()}


def _powerpoint_document_open_verification(file_path: Path, process_ids: list[int]) -> bool | None:
    if not process_ids:
        return False
    script = r'''
on run argv
    set targetDirectory to item 1 of argv
    set resolvedDirectory to item 2 of argv
    set targetName to item 3 of argv
    tell application "Microsoft PowerPoint"
        repeat with deckIndex from 1 to count of presentations
            try
                if (get name of presentation deckIndex) is (my targetName) then
                    if ((get path of presentation deckIndex) as text) is (my targetDirectory) or ((get path of presentation deckIndex) as text) is (my resolvedDirectory) then return "true"
                end if
            end try
        end repeat
    end tell
    return "false"
end run
'''
    raw_directory = str(file_path.parent)
    resolved_directory = str(file_path.parent.resolve())
    try:
        result = _run(["osascript", "-e", script, raw_directory, resolved_directory, file_path.name], check=False, timeout=8)
    except (OSError, subprocess.TimeoutExpired):
        return None
    if result.returncode != 0:
        return None
    value = result.stdout.strip().lower()
    if value == "true":
        return True
    if value == "false":
        return False
    return None


def _document_open_verification(file_path: Path, process_ids: list[int], host_name: str | None = None) -> bool | None:
    """Return a verified open state, or None when this platform cannot prove it."""
    if sys.platform == "darwin" and host_name == "powerpoint":
        # PowerPoint may release its file descriptor after loading a deck, so
        # lsof alone can falsely report that a visible presentation is closed.
        return _powerpoint_document_open_verification(file_path, process_ids)
    if sys.platform != "darwin" or not shutil.which("lsof"):
        return None
    opened = _run(["lsof", "-F", "p", "--", str(file_path.resolve())], check=False, timeout=15)
    return bool(_parse_lsof_process_ids(opened.stdout).intersection(process_ids))


def _wait_for_document_open(file_path: Path, host_name: str, host: dict) -> bool | None:
    if sys.platform != "darwin":
        return None
    timeout_ms = max(0, min(30000, int(os.environ.get("SCIENTIFIC_ILLUSTRATOR_OPEN_VERIFY_TIMEOUT_MS", "5000"))))
    deadline = time.monotonic() + timeout_ms / 1000
    while True:
        verified = _document_open_verification(file_path, _main_process_ids(host_name, host), host_name)
        if verified is True or verified is None or time.monotonic() >= deadline:
            return verified
        time.sleep(0.1)


def _dispatch_macos_open(file_path: Path, host: dict, focus_policy: str) -> dict:
    common = ["-g"] if focus_policy == "preserve" else []
    commands = []
    if host.get("bundle_id"):
        commands.append((["open", *common, "-b", host["bundle_id"], str(file_path)], "launch-services-bundle-id"))
    if host.get("path"):
        commands.append((["open", *common, "-a", host["path"], str(file_path)], "launch-services-app-path"))
    attempts = []
    for command, method in commands:
        try:
            dispatched = _run(command, check=False, timeout=30)
            attempts.append({"method": method, "return_code": dispatched.returncode, "stderr": dispatched.stderr.strip()})
        except (OSError, subprocess.TimeoutExpired) as exc:
            attempts.append({"method": method, "return_code": None, "stderr": str(exc)})
            continue
        if dispatched.returncode == 0:
            return {
                "open_dispatched": True,
                "dispatch_method": method,
                "dispatch_return_code": 0,
                "dispatch_attempts": attempts,
            }
    return {
        "open_dispatched": False,
        "dispatch_method": attempts[-1]["method"] if attempts else None,
        "dispatch_return_code": attempts[-1]["return_code"] if attempts else None,
        "dispatch_attempts": attempts,
        "dispatch_error": attempts[-1]["stderr"] if attempts else "No macOS application target was available.",
    }


def _presentation_host_info(args: dict | None = None, state: dict | None = None) -> dict:
    available = _available_hosts()
    host_name, host = _select_host(args, state, available)
    installed = bool(host["installed"])
    process_ids = _main_process_ids(host_name, host)
    presentation_count = None
    if process_ids and sys.platform == "darwin" and host_name == "powerpoint":
        count_text = _osascript('tell application "Microsoft PowerPoint" to get count of presentations')
        try:
            presentation_count = int(count_text)
        except (TypeError, ValueError):
            presentation_count = None
    return {
        "platform": sys.platform,
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": host_name,
        "target_application": host_name,
        "microsoft_powerpoint_used": host_name == "powerpoint",
        "available_hosts": available,
        "application_path": host.get("path"),
        "application_bundle_id": host.get("bundle_id"),
        "application_executable_path": host.get("executable_path"),
        "installed": installed,
        "wps_installed": bool(available["wps"]["installed"]),
        "application_version": host.get("version", ""),
        "running_processes": len(process_ids),
        "main_process_running": bool(process_ids),
        "process_ids": process_ids,
        "active_application_process_id": process_ids[0] if process_ids else 0,
        "presentation_count": presentation_count,
    }


def _refresh_presentation(file_path: Path, state: dict | None = None, *, focus_policy: str | None = None) -> dict:
    focus_policy = focus_policy or _focus_policy()
    host_name, host = _select_host(state=state)
    result = {
        "target_application": host_name,
        "microsoft_powerpoint_used": host_name == "powerpoint",
        "file_path": str(file_path),
        "focus_policy": focus_policy,
        "sync_enabled": os.environ.get("SCIENTIFIC_ILLUSTRATOR_POWERPOINT_SYNC", "1") != "0",
        "open_dispatched": False,
        "dispatch_method": None,
        "dispatch_return_code": None,
        "document_open_before": None,
        "document_open_verified": None,
        "refresh_verified": None,
    }
    if os.environ.get("SCIENTIFIC_ILLUSTRATOR_POWERPOINT_SYNC", "1") == "0":
        result["dispatch_error"] = "Application synchronization is disabled by SCIENTIFIC_ILLUSTRATOR_POWERPOINT_SYNC=0."
        if state is not None:
            state["last_refresh"] = result
            state["refresh_pending"] = True
            _write_state(state)
        return result
    if not host["installed"]:
        result["dispatch_error"] = f"{host_name} is not installed or its main executable could not be resolved."
        if state is not None:
            state["last_refresh"] = result
            state["refresh_pending"] = True
            _write_state(state)
        return result
    if not file_path.exists():
        result["dispatch_error"] = f"Managed presentation does not exist: {file_path}"
        if state is not None:
            state["last_refresh"] = result
            state["refresh_pending"] = True
            _write_state(state)
        return result

    process_ids = _main_process_ids(host_name, host)
    result["document_open_before"] = _document_open_verification(file_path, process_ids, host_name)
    if sys.platform == "win32":
        try:
            result.update(_open_windows_presentation(file_path, host.get("executable_path") or host.get("path"), focus_policy))
        except OSError as exc:
            result["dispatch_error"] = str(exc)
    elif sys.platform == "darwin" and host_name == "powerpoint" and process_ids:
        if result["document_open_before"] is not True:
            # Do not query every PowerPoint document when the exact managed
            # file is not open. An unrelated modal document can block
            # AppleScript even though there is nothing for us to reload.
            result.update(_dispatch_macos_open(file_path, host, focus_policy))
        else:
            close_script = r'''
on run argv
    set targetDirectory to item 1 of argv
    set resolvedDirectory to item 2 of argv
    set targetName to item 3 of argv
    tell application "Microsoft PowerPoint"
        repeat with deckIndex from 1 to count of presentations
            try
                if (get name of presentation deckIndex) is (my targetName) then
                    if ((get path of presentation deckIndex) as text) is (my targetDirectory) or ((get path of presentation deckIndex) as text) is (my resolvedDirectory) then
                        if (get saved of presentation deckIndex) then
                            close presentation deckIndex saving no
                            return "closed"
                        else
                            return "unsaved"
                        end if
                    end if
                end if
            end try
        end repeat
    end tell
    return "not-found"
end run
'''
            try:
                close_attempt = _run([
                    "osascript", "-e", close_script,
                    str(file_path.parent), str(file_path.parent.resolve()), file_path.name,
                ], check=False, timeout=8)
                result["existing_powerpoint_document"] = close_attempt.stdout.strip() or "unknown"
                if result["existing_powerpoint_document"] == "unsaved":
                    result["dispatch_error"] = "The managed PowerPoint document has unsaved user changes; refusing to close and reload it. Save or discard those edits explicitly, then retry refresh."
                    result["reload_blocked_by_unsaved_changes"] = True
                else:
                    result.update(_dispatch_macos_open(file_path, host, focus_policy))
            except subprocess.TimeoutExpired:
                # A modal/read-only dialog must never make an MCP edit hang.
                # The managed file is still safely updated on disk and can be
                # reopened after the user dismisses the dialog.
                result["dispatch_error"] = "PowerPoint did not answer the safe reload check within 8 seconds; no close/reopen was attempted."
    elif sys.platform == "darwin":
        result.update(_dispatch_macos_open(file_path, host, focus_policy))
    else:
        result["dispatch_error"] = f"Application refresh is unsupported on platform {sys.platform}."

    if result["open_dispatched"]:
        result["document_open_verified"] = _wait_for_document_open(file_path, host_name, host)
    if not result["open_dispatched"] or result["document_open_verified"] is False:
        result["refresh_verified"] = False
    elif sys.platform == "darwin" and host_name == "powerpoint":
        result["refresh_verified"] = result["document_open_verified"]
    elif sys.platform == "darwin" and host_name == "wps" and result["document_open_before"] is False:
        result["refresh_verified"] = result["document_open_verified"]
    else:
        result["refresh_verified"] = None
        result["verification_note"] = "The open request succeeded, but this application backend does not expose a reliable document-reload acknowledgement."
    result["request_succeeded"] = bool(result["open_dispatched"] and result["document_open_verified"] is not False)
    result["checked_at_unix"] = time.time()
    if state is not None:
        state["last_refresh"] = result
        # A dispatched request is not the same thing as a verified reload. In
        # particular, Windows WPS exposes no reliable acknowledgement and an
        # already-open Mac WPS document may only prove that the file is open.
        state["refresh_pending"] = result["refresh_verified"] is not True
        _write_state(state)
    return result


def _managed_path(label: str = "scientific-illustrator") -> Path:
    STATE_DIR.mkdir(parents=True, exist_ok=True)
    return STATE_DIR / f"{label}-{time.strftime('%Y%m%d-%H%M%S')}-{uuid.uuid4().hex[:6]}.pptx"


def _new_presentation(path: Path) -> Presentation:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return prs


def _require_path(create: bool = False) -> tuple[dict, Path]:
    state = _state()
    raw_path = state.get("path")
    if raw_path:
        path = Path(raw_path)
        if path.exists():
            return state, path
    if not create:
        raise RuntimeError("No managed PowerPoint presentation is active. Call powerpoint_launch or powerpoint_new_presentation first.")
    path = _managed_path()
    _new_presentation(path)
    host_name, _ = _select_host(state=state)
    state.update({"path": str(path), "source_path": None, "owned": True, "read_only": False, "metadata": {}, "host_application": host_name})
    _write_state(state)
    return state, path


def _load(writable: bool = True) -> tuple[dict, Path, Presentation]:
    state, path = _require_path(False)
    if writable and state.get("read_only") is True:
        raise PermissionError("The managed presentation was launched read_only=true. Reopen it as a working copy before editing.")
    return state, path, Presentation(path)


def _save(prs: Presentation, state: dict, path: Path, *, refresh: bool = True) -> None:
    temp_path = path.with_suffix(".saving.pptx")
    prs.save(temp_path)
    os.replace(temp_path, path)
    state["managed_file_mtime_ns"] = path.stat().st_mtime_ns
    state["refresh_pending"] = True
    _write_state(state)
    deferred = os.environ.get("SCIENTIFIC_ILLUSTRATOR_DEFER_REFRESH", "0") == "1"
    if refresh and not deferred:
        _refresh_presentation(path, state)


def _slide(prs: Presentation, index: int):
    if index < 1 or index > len(prs.slides):
        raise ValueError(f"slide_index {index} is outside 1..{len(prs.slides)}")
    return prs.slides[index - 1]


def _shape(slide, args: dict):
    wanted_name = args.get("shape_name")
    wanted_id = args.get("shape_id")
    matches = [
        shape for shape in slide.shapes
        if (wanted_name is not None and shape.name.lower() == str(wanted_name).lower())
        or (wanted_id is not None and shape.shape_id == int(wanted_id))
    ]
    if len(matches) > 1:
        raise ValueError(f"Shape target is ambiguous because the semantic name is duplicated: {wanted_name}")
    if matches:
        return matches[0]
    raise ValueError(f"Shape not found: {wanted_name or wanted_id}")


def _assert_shape_name_available(slide, name: str | None, exclude_shape_id: int | None = None) -> None:
    if name is None or not str(name).strip():
        return
    wanted = str(name).strip().lower()
    for existing in slide.shapes:
        if existing.shape_id != exclude_shape_id and existing.name.strip().lower() == wanted:
            raise ValueError(f"Shape name already exists on this slide: {name}")


def _pt(value: float | int) -> Pt:
    return Pt(float(value))


def _rgb(value: str | None) -> RGBColor | None:
    if not value:
        return None
    text = str(value).strip().lstrip("#")
    if len(text) != 6:
        raise ValueError(f"Invalid color {value}; expected #RRGGBB")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _set_alpha(color_element, transparency: float | int | None) -> None:
    if transparency is None or color_element is None:
        return
    value = max(0.0, min(100.0, float(transparency)))
    for child in list(color_element):
        if child.tag.endswith("}alpha"):
            color_element.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(round((100.0 - value) * 1000))))
    color_element.append(alpha)


def _set_fill(shape, color: str | None, transparency: float | int | None = None) -> None:
    if color is None and transparency is None:
        return
    if not hasattr(shape, "fill"):
        raise ValueError(f"Shape '{getattr(shape, 'name', '<unnamed>')}' does not support fill styling.")
    rgb = _rgb(color)
    if rgb is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    solid_fill = shape.fill._xPr.solidFill
    if transparency is not None and (solid_fill is None or len(solid_fill) == 0):
        raise ValueError("fill_transparency requires an explicit solid fill_color for this OOXML object.")
    _set_alpha(solid_fill[0] if solid_fill is not None and len(solid_fill) else None, transparency)


def _set_line(shape, args: dict) -> None:
    requested = any(args.get(key) is not None for key in ("line_color", "line_width", "line_dash", "line_transparency"))
    if not requested:
        return
    if not hasattr(shape, "line"):
        raise ValueError(f"Shape '{getattr(shape, 'name', '<unnamed>')}' does not support line styling.")
    rgb = _rgb(args.get("line_color"))
    if rgb is not None:
        shape.line.color.rgb = rgb
    if args.get("line_width") is not None:
        shape.line.width = _pt(args["line_width"])
    dash = args.get("line_dash")
    if dash is not None:
        dash_styles = {
            "solid": MSO_LINE_DASH_STYLE.SOLID,
            "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
            "round_dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
            "dash": MSO_LINE_DASH_STYLE.DASH,
            "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
            "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
            "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
        }
        if str(dash) == "long_dash_dot_dot":
            line = shape.line._get_or_add_ln()
            preset = line.find("a:prstDash", line.nsmap)
            if preset is None:
                preset = OxmlElement("a:prstDash")
                arrow = line.find("a:headEnd", line.nsmap)
                if arrow is None:
                    arrow = line.find("a:tailEnd", line.nsmap)
                if arrow is None:
                    line.append(preset)
                else:
                    arrow.addprevious(preset)
            preset.set("val", "lgDashDotDot")
        else:
            shape.line.dash_style = dash_styles[str(dash)]
    if args.get("line_transparency") is not None:
        line = shape.line._get_or_add_ln()
        solid_fill = line.solidFill
        if solid_fill is None or len(solid_fill) == 0:
            raise ValueError("line_transparency requires an explicit solid line_color for this OOXML object.")
        _set_alpha(solid_fill[0] if solid_fill is not None and len(solid_fill) else None, args["line_transparency"])


def _set_arrow(shape, start_arrow: str = "none", end_arrow: str = "none") -> None:
    if not hasattr(shape, "line"):
        return
    line = shape.line._get_or_add_ln()
    for tag, kind in (("a:headEnd", start_arrow), ("a:tailEnd", end_arrow)):
        existing = line.find(tag, line.nsmap)
        if existing is None:
            existing = OxmlElement(tag)
            line.append(existing)
        mapping = {"none": "none", "open": "arrow", "triangle": "triangle", "stealth": "stealth", "diamond": "diamond", "oval": "oval"}
        existing.set("type", mapping.get(str(kind), "none"))


def _get_arrow(shape, tag: str) -> str:
    if not hasattr(shape, "line"):
        return "none"
    line = shape.line._get_or_add_ln()
    element = line.find(tag, line.nsmap)
    inverse = {"none": "none", "arrow": "open", "triangle": "triangle", "stealth": "stealth", "diamond": "diamond", "oval": "oval"}
    return inverse.get(element.get("type") if element is not None else "none", "none")


def _apply_text_frame(text_frame, args: dict) -> None:
    if args.get("margin_left") is not None:
        text_frame.margin_left = _pt(args["margin_left"])
    if args.get("margin_right") is not None:
        text_frame.margin_right = _pt(args["margin_right"])
    if args.get("margin_top") is not None:
        text_frame.margin_top = _pt(args["margin_top"])
    if args.get("margin_bottom") is not None:
        text_frame.margin_bottom = _pt(args["margin_bottom"])
    if args.get("word_wrap") is not None:
        text_frame.word_wrap = bool(args["word_wrap"])
    auto = args.get("text_autofit")
    if auto == "shrink_text":
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif auto == "grow_shape":
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    elif auto == "none":
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    valign = {"top": MSO_VERTICAL_ANCHOR.TOP, "middle": MSO_VERTICAL_ANCHOR.MIDDLE, "bottom": MSO_VERTICAL_ANCHOR.BOTTOM}
    if args.get("vertical_alignment") in valign:
        text_frame.vertical_anchor = valign[args["vertical_alignment"]]
    alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    font_color = _rgb(args.get("font_color"))
    for paragraph in text_frame.paragraphs:
        if args.get("alignment") in alignment:
            paragraph.alignment = alignment[args["alignment"]]
        if not paragraph.runs:
            paragraph.add_run()
        for run in paragraph.runs:
            font = run.font
            if args.get("font_name") is not None:
                font.name = str(args["font_name"])
            if args.get("font_size") is not None:
                font.size = _pt(args["font_size"])
            if args.get("bold") is not None:
                font.bold = bool(args["bold"])
            if args.get("italic") is not None:
                font.italic = bool(args["italic"])
            if font_color is not None:
                font.color.rgb = font_color


def _set_text(shape, text: object, args: dict) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.text = "" if text is None else str(text)
    _apply_text_frame(shape.text_frame, args)


SHAPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "round_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "hexagon": MSO_SHAPE.HEXAGON,
    "pentagon": MSO_SHAPE.PENTAGON,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "chevron": MSO_SHAPE.CHEVRON,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "cloud": MSO_SHAPE.CLOUD,
    "arc": MSO_SHAPE.ARC,
    "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
}

CHARTS = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


def _shape_result(shape, slide_index: int) -> dict:
    return {
        "slide_index": slide_index,
        "shape_id": shape.shape_id,
        "shape_name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt,
    }


def action_status(args: dict) -> dict:
    state = _state()
    host = _presentation_host_info(args, state)
    path = Path(state["path"]) if state.get("path") else None
    managed_file_exists = bool(path and path.exists())
    document_open_verified = _document_open_verification(path, host["process_ids"], host["host_application"]) if managed_file_exists else False
    last_refresh = state.get("last_refresh") or {}
    return {
        **host,
        "connected_to_active_application": False,
        "connection_mode": "file-backed-working-copy",
        "active_presentation": document_open_verified is True,
        "managed_file_exists": managed_file_exists,
        "document_open_verified": document_open_verified,
        "document_open_verification_available": document_open_verified is not None,
        "open_dispatched": bool(last_refresh.get("open_dispatched", False)),
        "refresh_verified": last_refresh.get("refresh_verified"),
        "refresh_pending": bool(state.get("refresh_pending", False)),
        "last_refresh": state.get("last_refresh"),
        "managed_path": str(path) if path else None,
        "source_path": state.get("source_path"),
        "read_only": bool(state.get("read_only", False)),
        "native_editable_output": True,
        "focus_policy": _focus_policy(),
        "live_semantics": "editable file-backed OOXML; object mutations update the managed PPTX, and application refresh is dispatched only at explicit or sequence checkpoints",
    }


def action_capabilities(args: dict) -> dict:
    host = action_status(args)
    supported = {
        "text_box": "Shapes.add_textbox",
        "auto_shape": "Shapes.add_shape",
        "free_line_or_arrow": "Shapes.add_connector",
        "attached_connector": "file-backed native connector",
        "table": "Shapes.add_table",
        "chart": "Shapes.add_chart",
        "picture_or_svg": "Shapes.add_picture",
        "duplicate": "OOXML duplicate",
        "group": "Shapes.add_group_shape",
        "ungroup": "OOXML ungroup",
        "z_order": "OOXML tree order",
        "align": "exact geometry",
        "distribute": "exact geometry",
        "figure_audit": "deterministic OOXML geometry audit",
    }
    return {
        "detection": {
            "read_only": True,
            "launched_powerpoint": False,
            "active_deck_modified": False,
            "basis": ["python-pptx OOXML API", "application bundle or executable discovery", "exact main-process matching", "open-file verification where the operating system exposes it"],
        },
        "host": host,
        "native_object_families": [
            {"family": key, "powerpoint_api": value, "host_supported": True, "editable": key != "figure_audit"}
            for key, value in supported.items()
        ],
        "auto_shapes": [{"plugin_name": name, "value": int(value)} for name, value in SHAPES.items()] if args.get("include_auto_shapes", True) else [],
        "chart_types": [{"plugin_name": name, "value": int(value)} for name, value in CHARTS.items()] if args.get("include_chart_types", True) else [],
        "connector_types": [{"plugin_name": name} for name in ("straight", "elbow", "curve")],
        "arrowhead_styles": [{"plugin_name": name} for name in ("none", "open", "triangle", "stealth", "diamond", "oval")],
        "limitations": [
            "PowerPoint for Mac and WPS Presentation use an isolated file-backed working copy because they do not expose the Windows PowerPoint COM automation server used by this plugin.",
            "File-backed drawing cannot claim in-memory attachment to an arbitrary current WPS document. The managed working copy is authoritative and application refresh is checkpointed to avoid repeated focus and reload requests.",
            "WPS does not expose a reliable reload acknowledgement. Status distinguishes an accepted open request, a verified open file, and an unverified refresh instead of reporting false success.",
            "Renderer exports use local LibreOffice/Poppler when available and remain separate from the editable PPTX.",
        ],
    }


def action_launch(args: dict) -> dict:
    state = _state()
    host_name, _ = _select_host(args, state)
    state["host_application"] = host_name
    supplied = args.get("file_path")
    if supplied:
        source = Path(supplied).expanduser().resolve()
        if not source.exists():
            raise FileNotFoundError(source)
        if source.suffix.lower() not in {".pptx", ".pptm", ".ppsx"}:
            raise ValueError("file_path must end with .pptx, .pptm, or .ppsx.")
        if source.suffix.lower() != ".pptx" and not args.get("read_only", False):
            raise ValueError("The editable OOXML working-copy backend accepts only .pptx. Editing .pptm or .ppsx could discard macros or change presentation semantics; save a PPTX copy first or use Windows PowerPoint COM.")
        if args.get("read_only", False):
            working = source
            owned = False
        else:
            working = _managed_path(source.stem + "-working")
            shutil.copy2(source, working)
            owned = True
        state.update({"path": str(working), "source_path": str(source), "owned": owned, "read_only": bool(args.get("read_only", False)), "metadata": {}, "last_refresh": None, "refresh_pending": True})
    elif state.get("path") and Path(state["path"]).exists():
        working = Path(state["path"])
    elif args.get("create_if_missing", True):
        working = _managed_path()
        _new_presentation(working)
        state.update({"path": str(working), "source_path": None, "owned": True, "read_only": False, "metadata": {}, "last_refresh": None, "refresh_pending": True})
    else:
        raise RuntimeError("No managed presentation exists and create_if_missing=false.")
    _write_state(state)
    refresh_result = None
    if args.get("visible", True):
        refresh_result = _refresh_presentation(working, state)
    return {
        **action_status({"host_application": host_name}),
        "opened_path": str(working),
        "working_copy": bool(state.get("owned")),
        "open_result": refresh_result,
    }


def action_new_presentation(args: dict) -> dict:
    path = _managed_path()
    prs = _new_presentation(path)
    host_name, _ = _select_host(args)
    state = {"path": str(path), "source_path": None, "owned": True, "read_only": False, "metadata": {}, "host_application": host_name, "last_refresh": None, "refresh_pending": True}
    _write_state(state)
    refresh_result = _refresh_presentation(path, state)
    return {
        "path": str(path),
        "slide_count": len(prs.slides),
        "platform": sys.platform,
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": host_name,
        "target_application": host_name,
        "microsoft_powerpoint_used": host_name == "powerpoint",
        "open_result": refresh_result,
        "active_presentation": refresh_result["document_open_verified"] is True,
    }


def _inventory(prs: Presentation, args: dict) -> list[dict]:
    result = []
    max_slides = int(args.get("max_slides", 100))
    max_shapes = int(args.get("max_shapes_per_slide", 200))
    for index, slide in enumerate(list(prs.slides)[:max_slides], 1):
        shapes = []
        for shape in list(slide.shapes)[:max_shapes]:
            item = _shape_result(shape, index)
            if args.get("include_text", True) and getattr(shape, "has_text_frame", False):
                item["text"] = shape.text
            item["is_picture"] = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            item["is_table"] = bool(getattr(shape, "has_table", False))
            item["is_chart"] = bool(getattr(shape, "has_chart", False))
            shapes.append(item)
        result.append({"slide_index": index, "shape_count": len(slide.shapes), "shapes": shapes})
    return result


def action_inspect(args: dict) -> dict:
    state, path, prs = _load(False)
    return {
        "path": str(path),
        "platform": sys.platform,
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": state.get("host_application", "auto"),
        "slide_width": prs.slide_width.pt,
        "slide_height": prs.slide_height.pt,
        "slide_count": len(prs.slides),
        "slides": _inventory(prs, args),
    }


def action_add_slide(args: dict) -> dict:
    state, path, prs = _load()
    layout_index = {"blank": 6, "title": 0, "text": 1}.get(args.get("layout", "blank"), 6)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    position = args.get("position")
    if position and int(position) <= len(prs.slides):
        slide_ids = prs.slides._sldIdLst
        slide_id = slide_ids[-1]
        slide_ids.remove(slide_id)
        slide_ids.insert(max(0, int(position) - 1), slide_id)
    index = list(prs.slides).index(slide) + 1
    if args.get("name"):
        state.setdefault("metadata", {}).setdefault("slides", {})[str(index)] = {"name": args["name"]}
    _save(prs, state, path)
    return {"slide_index": index, "slide_count": len(prs.slides), "name": args.get("name")}


def action_activate_slide(args: dict) -> dict:
    state, path = _require_path(False)
    index = int(args["slide_index"])
    prs = Presentation(path)
    _slide(prs, index)
    refresh_result = _refresh_presentation(path, state, focus_policy="foreground")
    activated = refresh_result["document_open_verified"] is True
    return {
        "slide_index": index,
        "activated": activated,
        "document_activation_verified": activated,
        "exact_slide_selection_verified": None,
        "open_result": refresh_result,
        "note": "A foreground open request was dispatched. activated=true only when the operating system verified the managed file as open; exact slide selection is not exposed by the WPS/OOXML backend.",
    }


def action_refresh(args: dict) -> dict:
    state, path = _require_path(False)
    policy = args.get("focus_policy") or _focus_policy()
    return _refresh_presentation(path, state, focus_policy=policy)


def action_add_textbox(args: dict) -> dict:
    state, path, prs = _load()
    slide = _slide(prs, int(args["slide_index"]))
    _assert_shape_name_available(slide, args.get("name"))
    shape = slide.shapes.add_textbox(_pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    _set_text(shape, args.get("text", ""), args)
    _set_fill(shape, args.get("fill_color"), args.get("fill_transparency"))
    _set_line(shape, args)
    if args.get("start_arrow") is not None or args.get("end_arrow") is not None:
        _set_arrow(
            shape,
            args.get("start_arrow", _get_arrow(shape, "a:headEnd")),
            args.get("end_arrow", _get_arrow(shape, "a:tailEnd")),
        )
    _save(prs, state, path)
    return _shape_result(shape, int(args["slide_index"]))


def _normalize_enum_name(value: object, prefixes: tuple[str, ...] = ()) -> str:
    text = str(value).strip()
    for prefix in prefixes:
        if text.lower().startswith(prefix.lower()):
            text = text[len(prefix):]
            break
    text = re.sub(r"([a-z0-9])([A-Z])", r"\1_\2", text)
    return re.sub(r"[^A-Za-z0-9]+", "_", text).strip("_").lower()


def _shape_enum(args: dict):
    if args.get("shape_type_id") is not None:
        return MSO_SHAPE(int(args["shape_type_id"]))
    key = _normalize_enum_name(args.get("shape", "rectangle"), ("msoShape",))
    if key not in SHAPES:
        raise ValueError(f"Unsupported OOXML auto shape '{args.get('shape')}'. Use one of: {', '.join(sorted(SHAPES))}")
    return SHAPES[key]


def action_add_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    shape = slide.shapes.add_shape(_shape_enum(args), _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    if args.get("rotation") is not None:
        shape.rotation = float(args["rotation"])
    _set_fill(shape, args.get("fill_color"), args.get("fill_transparency"))
    _set_line(shape, args)
    if args.get("text") is not None:
        _set_text(shape, args["text"], args)
    _save(prs, state, path)
    return _shape_result(shape, index)


def action_add_image(args: dict) -> dict:
    state, path, prs = _load()
    if args.get("atomic_raster_unit") is not True or args.get("contains_reconstructable_content") is not False:
        raise ValueError("OOXML pictures must be atomic_raster_unit=true and contains_reconstructable_content=false.")
    raster_reason = str(args.get("raster_reason", "")).strip()
    decomposition_note = str(args.get("decomposition_note", "")).strip()
    if len(raster_reason) < 8:
        raise ValueError("raster_reason must specifically explain why this atomic visual field cannot be rebuilt natively.")
    if len(decomposition_note) < 8:
        raise ValueError("decomposition_note must explain what was separated or why no finer split is possible.")
    crop_keys = [f"crop_{side}_{unit}" for side in ("left", "top", "right", "bottom") for unit in ("percent", "points")]
    has_crop = any(args.get(key) is not None for key in crop_keys)
    if args.get("source_is_tightly_cropped") is not True and not has_crop:
        raise ValueError("source_is_tightly_cropped=false requires at least one crop_* value.")
    image_path = Path(args["image_path"]).expanduser().resolve()
    if not image_path.exists():
        raise FileNotFoundError(image_path)
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    shape = slide.shapes.add_picture(str(image_path), _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    crop_values = {}
    for side in ("left", "top", "right", "bottom"):
        percent = args.get(f"crop_{side}_percent")
        points = args.get(f"crop_{side}_points")
        if percent is not None and points is not None:
            raise ValueError(f"Specify only one of crop_{side}_percent or crop_{side}_points.")
        if percent is not None:
            crop_values[side] = float(percent) / 100.0
        elif points is not None:
            base = float(args["width"] if side in ("left", "right") else args["height"])
            crop_values[side] = float(points) / base
        else:
            crop_values[side] = 0.0
    if crop_values["left"] + crop_values["right"] >= 1 or crop_values["top"] + crop_values["bottom"] >= 1:
        raise ValueError("Picture crop values remove the entire image width or height.")
    for side, value in crop_values.items():
        if value < 0 or value >= 1:
            raise ValueError(f"crop_{side} must be between 0 and the rendered picture dimension.")
        setattr(shape, f"crop_{side}", value)

    c_nv_pr = shape._element.nvPicPr.cNvPr
    c_nv_pr.set("descr", str(args.get("alt_text") or f"Raster-only visual evidence: {raster_reason}"))
    c_nv_pr.set("title", shape.name)
    c_nv_pic_pr = shape._element.nvPicPr.find("p:cNvPicPr", shape._element.nsmap)
    pic_locks = c_nv_pic_pr.find("a:picLocks", c_nv_pic_pr.nsmap) if c_nv_pic_pr is not None else None
    if pic_locks is not None:
        pic_locks.set("noChangeAspect", "1" if args.get("lock_aspect_ratio", False) else "0")
    state.setdefault("metadata", {}).setdefault("rasters", {})[str(shape.shape_id)] = {
        "slide_index": index,
        "shape_name": shape.name,
        "raster_reason": raster_reason,
        "atomic_raster_unit": True,
        "contains_reconstructable_content": False,
        "decomposition_note": decomposition_note,
        "source_is_tightly_cropped": bool(args.get("source_is_tightly_cropped")),
        "crop": crop_values,
    }
    _save(prs, state, path)
    return {**_shape_result(shape, index), "raster_declaration": state["metadata"]["rasters"][str(shape.shape_id)]}


def _trim_line(x1: float, y1: float, x2: float, y2: float, start: float, end: float) -> tuple[float, float, float, float]:
    dx, dy = x2 - x1, y2 - y1
    length = math.hypot(dx, dy)
    if length <= start + end or length == 0:
        raise ValueError("Line clearances consume the entire segment.")
    ux, uy = dx / length, dy / length
    return x1 + ux * start, y1 + uy * start, x2 - ux * end, y2 - uy * end


def _connector_endpoints(shape) -> tuple[float, float, float, float]:
    transform = shape._element.spPr.xfrm
    flip_h = str(transform.get("flipH", "0")).lower() in {"1", "true"}
    flip_v = str(transform.get("flipV", "0")).lower() in {"1", "true"}
    left, top, right, bottom = shape.left.pt, shape.top.pt, shape.left.pt + shape.width.pt, shape.top.pt + shape.height.pt
    return (
        right if flip_h else left,
        bottom if flip_v else top,
        left if flip_h else right,
        top if flip_v else bottom,
    )


def action_add_line(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    x1, y1, x2, y2 = _trim_line(
        float(args["begin_x"]), float(args["begin_y"]), float(args["end_x"]), float(args["end_y"]),
        float(args.get("start_clearance", 0)), float(args.get("end_clearance", 0)),
    )
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _pt(x1), _pt(y1), _pt(x2), _pt(y2))
    if args.get("name"):
        shape.name = args["name"]
    _set_line(shape, args)
    _set_arrow(shape, args.get("start_arrow", "none"), args.get("end_arrow", "none"))
    state.setdefault("metadata", {}).setdefault("lines", {})[str(shape.shape_id)] = {
        "slide_index": index,
        "shape_name": shape.name,
        "begin_x": x1,
        "begin_y": y1,
        "end_x": x2,
        "end_y": y2,
        "start_arrow": args.get("start_arrow", "none"),
        "end_arrow": args.get("end_arrow", "none"),
    }
    _save(prs, state, path)
    return {**_shape_result(shape, index), "begin_x": x1, "begin_y": y1, "end_x": x2, "end_y": y2}


def _site_point(shape, site: int) -> tuple[float, float]:
    left, top, width, height = shape.left.pt, shape.top.pt, shape.width.pt, shape.height.pt
    points = {
        1: (left + width / 2, top),
        2: (left + width, top + height / 2),
        3: (left + width / 2, top + height),
        4: (left, top + height / 2),
    }
    if int(site) not in points:
        raise ValueError("The OOXML backend supports connection sites 1..4: top, right, bottom, and left.")
    return points[int(site)]


def _automatic_connection_sites(source, target) -> tuple[int, int]:
    source_center = (source.left.pt + source.width.pt / 2, source.top.pt + source.height.pt / 2)
    target_center = (target.left.pt + target.width.pt / 2, target.top.pt + target.height.pt / 2)
    dx, dy = target_center[0] - source_center[0], target_center[1] - source_center[1]
    if abs(dx) >= abs(dy):
        return (2, 4) if dx >= 0 else (4, 2)
    return (3, 1) if dy >= 0 else (1, 3)


def _bind_connector(shape, source, source_site: int, target, target_site: int) -> None:
    properties = shape._element.nvCxnSpPr.cNvCxnSpPr
    for tag in ("a:stCxn", "a:endCxn"):
        existing = properties.find(tag, properties.nsmap)
        if existing is not None:
            properties.remove(existing)
    start = OxmlElement("a:stCxn")
    start.set("id", str(source.shape_id))
    start.set("idx", str(source_site - 1))
    end = OxmlElement("a:endCxn")
    end.set("id", str(target.shape_id))
    end.set("idx", str(target_site - 1))
    properties.append(start)
    properties.append(end)


def action_add_connector(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    source = _shape(slide, {"shape_name": args["source_name"]})
    target = _shape(slide, {"shape_name": args["target_name"]})
    automatic_source_site, automatic_target_site = _automatic_connection_sites(source, target)
    source_site = int(args.get("source_site", automatic_source_site))
    target_site = int(args.get("target_site", automatic_target_site))
    x1, y1 = _site_point(source, source_site)
    x2, y2 = _site_point(target, target_site)
    connector_type = {"straight": MSO_CONNECTOR.STRAIGHT, "elbow": MSO_CONNECTOR.ELBOW, "curve": MSO_CONNECTOR.CURVE}[args.get("connector_type", "elbow")]
    shape = slide.shapes.add_connector(connector_type, _pt(x1), _pt(y1), _pt(x2), _pt(y2))
    if args.get("name"):
        shape.name = args["name"]
    _set_line(shape, args)
    _set_arrow(shape, args.get("start_arrow", "none"), args.get("end_arrow", "triangle"))
    _bind_connector(shape, source, source_site, target, target_site)
    state.setdefault("metadata", {}).setdefault("connectors", {})[str(shape.shape_id)] = {
        "slide_index": index,
        "shape_name": shape.name,
        "source_name": source.name,
        "source_id": source.shape_id,
        "source_site": source_site,
        "target_name": target.name,
        "target_id": target.shape_id,
        "target_site": target_site,
        "connector_type": args.get("connector_type", "elbow"),
        "start_arrow": args.get("start_arrow", "none"),
        "end_arrow": args.get("end_arrow", "triangle"),
    }
    _save(prs, state, path)
    return {
        **_shape_result(shape, index),
        "source_name": source.name,
        "source_site": source_site,
        "target_name": target.name,
        "target_site": target_site,
        "attachment_mode": "ooxml-connection-site",
    }


def _set_cell_borders(cell, args: dict) -> None:
    color = _rgb(args.get("border_color"))
    width = args.get("border_width")
    if color is None and width is None:
        return
    color = color or RGBColor(0, 0, 0)
    width_emu = str(int(_pt(width if width is not None else 1)))
    tc_pr = cell._tc.get_or_add_tcPr()
    for index, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
        existing = tc_pr.find(tag, tc_pr.nsmap)
        if existing is not None:
            tc_pr.remove(existing)
        line = OxmlElement(tag)
        line.set("w", width_emu)
        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", str(color))
        solid_fill.append(srgb)
        line.append(solid_fill)
        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        line.append(dash)
        line.append(OxmlElement("a:round"))
        tc_pr.insert(index, line)


def _style_cell(cell, args: dict) -> None:
    if args.get("text") is not None:
        cell.text = str(args["text"])
    _set_fill(cell, args.get("fill_color"), args.get("fill_transparency"))
    _apply_text_frame(cell.text_frame, args)
    _set_cell_borders(cell, args)
    margin = args.get("cell_margin")
    if margin is not None:
        cell.margin_left = cell.margin_right = cell.margin_top = cell.margin_bottom = _pt(margin)


def action_add_table(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    rows, columns = int(args["rows"]), int(args["columns"])
    data = args.get("data") or []
    if len(data) > rows:
        raise ValueError(f"Table data has {len(data)} rows but rows={rows}.")
    for row_index, values in enumerate(data, 1):
        if len(values) > columns:
            raise ValueError(f"Table data row {row_index} has more than columns={columns} values.")
    for override in args.get("cell_styles") or []:
        row, col = int(override["row"]), int(override["column"])
        if row < 1 or row > rows or col < 1 or col > columns:
            raise ValueError(f"cell_styles entry ({row},{col}) is outside the table bounds {rows} x {columns}.")
    shape = slide.shapes.add_table(rows, columns, _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]))
    if args.get("name"):
        shape.name = args["name"]
    table = shape.table
    header_rows = min(rows, int(args.get("header_rows", 1)))
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            cell_args = dict(args)
            if row < len(data) and col < len(data[row]):
                cell_args["text"] = "" if data[row][col] is None else data[row][col]
            if row < header_rows:
                cell_args["fill_color"] = args.get("header_fill_color", args.get("fill_color"))
                cell_args["font_color"] = args.get("header_font_color", args.get("font_color"))
                cell_args["bold"] = args.get("header_bold", True)
            elif args.get("banded_rows") and (row - header_rows) % 2 == 0:
                cell_args["fill_color"] = args.get("band_fill_color", args.get("fill_color"))
            _style_cell(table.cell(row, col), cell_args)
    for override in args.get("cell_styles") or []:
        row, col = int(override["row"]) - 1, int(override["column"]) - 1
        _style_cell(table.cell(row, col), {**args, **override})
    _save(prs, state, path)
    return {**_shape_result(shape, index), "rows": len(table.rows), "columns": len(table.columns)}


def action_update_table_cell(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    if not getattr(shape, "has_table", False):
        raise ValueError(f"{shape.name} is not a table")
    row, col = int(args["row"]) - 1, int(args["column"]) - 1
    _style_cell(shape.table.cell(row, col), args)
    _save(prs, state, path)
    return {**_shape_result(shape, index), "row": row + 1, "column": col + 1, "text": shape.table.cell(row, col).text}


def action_update_table_layout(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    if not getattr(shape, "has_table", False):
        raise ValueError(f"{shape.name} is not a table")
    column_widths = args.get("column_widths")
    row_heights = args.get("row_heights")
    if column_widths is None and row_heights is None:
        raise ValueError("Provide column_widths and/or row_heights.")
    if column_widths is not None and len(column_widths) != len(shape.table.columns):
        raise ValueError(f"column_widths count {len(column_widths)} does not match table column count {len(shape.table.columns)}.")
    if row_heights is not None and len(row_heights) != len(shape.table.rows):
        raise ValueError(f"row_heights count {len(row_heights)} does not match table row count {len(shape.table.rows)}.")
    for column, width in zip(shape.table.columns, column_widths or []):
        column.width = _pt(width)
    for row, height in zip(shape.table.rows, row_heights or []):
        row.height = _pt(height)
    _save(prs, state, path)
    return {**_shape_result(shape, index), "column_widths": column_widths, "row_heights": row_heights}


def action_add_chart(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    categories = list(args["categories"])
    for series in args["series"]:
        if len(series["values"]) != len(categories):
            raise ValueError(f"Chart series '{series['name']}' has {len(series['values'])} values for {len(categories)} categories.")
    if args.get("chart_type_id") is not None:
        chart_type = XL_CHART_TYPE(int(args["chart_type_id"]))
        chart_key = next((name for name, value in CHARTS.items() if value == chart_type), str(chart_type))
    else:
        chart_key = _normalize_enum_name(args.get("chart_type", "column_clustered"), ("xl",))
        if chart_key not in CHARTS:
            raise ValueError(f"Unsupported OOXML chart type '{args.get('chart_type')}'. Use one of: {', '.join(sorted(CHARTS))}")
        chart_type = CHARTS[chart_key]
    if chart_key == "scatter":
        from pptx.chart.data import XyChartData
        data = XyChartData()
        try:
            x_values = [float(value) for value in categories]
        except (TypeError, ValueError) as exc:
            raise ValueError("Scatter-chart categories must be numeric x values.") from exc
        for series in args["series"]:
            data_series = data.add_series(str(series["name"]))
            for x_value, y_value in zip(x_values, series["values"]):
                data_series.add_data_point(x_value, float(y_value))
    else:
        data = CategoryChartData()
        data.categories = categories
        for series in args["series"]:
            data.add_series(str(series["name"]), tuple(series["values"]))
    shape = slide.shapes.add_chart(chart_type, _pt(args["left"]), _pt(args["top"]), _pt(args["width"]), _pt(args["height"]), data)
    if args.get("name"):
        shape.name = args["name"]
    chart = shape.chart
    if args.get("title") is not None:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(args["title"])
    chart.has_legend = bool(args.get("has_legend", True))
    if chart.has_legend:
        chart.legend.position = {
            "right": XL_LEGEND_POSITION.RIGHT,
            "left": XL_LEGEND_POSITION.LEFT,
            "top": XL_LEGEND_POSITION.TOP,
            "bottom": XL_LEGEND_POSITION.BOTTOM,
        }[args.get("legend_position", "right")]
    if args.get("chart_style") is not None:
        chart.chart_style = int(args["chart_style"])
    for argument_name, axis_name in (("category_axis_title", "category_axis"), ("value_axis_title", "value_axis")):
        if args.get(argument_name) is None:
            continue
        try:
            axis = getattr(chart, axis_name)
        except ValueError as exc:
            raise ValueError(f"Chart type '{chart_key}' does not expose {axis_name.replace('_', ' ')} for {argument_name}.") from exc
        axis.has_title = True
        axis.axis_title.text_frame.text = str(args[argument_name])
    _save(prs, state, path)
    return {**_shape_result(shape, index), "chart_type": str(chart_type), "series_count": len(args["series"])}


def _contains_chart(shape) -> bool:
    if bool(getattr(shape, "has_chart", False)):
        return True
    if shape_type_name(shape) == "GROUP":
        return any(_contains_chart(member) for member in shape.shapes)
    return False


def _remap_cloned_shape_ids(slide, new_element, root_name: str) -> tuple[dict[int, int], dict[int, str]]:
    existing_ids = [int(item.get("id")) for item in slide._element.xpath(".//p:cNvPr") if str(item.get("id", "")).isdigit()]
    next_id = max(existing_ids, default=1) + 1
    id_mapping: dict[int, int] = {}
    names_by_new_id: dict[int, str] = {}
    properties = list(new_element.xpath(".//p:cNvPr"))
    for index, item in enumerate(properties):
        old_id = int(item.get("id"))
        new_id = next_id + index
        old_name = str(item.get("name") or f"shape-{old_id}")
        new_name = root_name if index == 0 else f"{root_name}__{old_name}"
        item.set("id", str(new_id))
        item.set("name", new_name[:255])
        id_mapping[old_id] = new_id
        names_by_new_id[new_id] = new_name[:255]
    for connection in new_element.xpath(".//a:stCxn | .//a:endCxn"):
        old_id = int(connection.get("id"))
        if old_id in id_mapping:
            connection.set("id", str(id_mapping[old_id]))
    return id_mapping, names_by_new_id


def action_duplicate_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    source = _shape(slide, args)
    _assert_shape_name_available(slide, args["new_name"])
    if _contains_chart(source):
        raise ValueError("OOXML chart duplication is refused because copied chart shapes would share one embedded data part. Recreate the chart with powerpoint_add_chart and the original series instead.")
    new_element = copy.deepcopy(source._element)
    id_mapping, names_by_new_id = _remap_cloned_shape_ids(slide, new_element, args["new_name"])
    slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    duplicate = list(slide.shapes)[-1]
    duplicate.name = args["new_name"]
    for prop in ("left", "top", "width", "height"):
        if args.get(prop) is not None:
            setattr(duplicate, prop, _pt(args[prop]))
    if args.get("rotation") is not None:
        duplicate.rotation = float(args["rotation"])
    metadata = state.setdefault("metadata", {})
    for collection in ("rasters", "lines", "connectors"):
        collection_metadata = metadata.get(collection, {})
        for old_id, new_id in id_mapping.items():
            source_metadata = collection_metadata.get(str(old_id))
            if source_metadata is None:
                continue
            copied_metadata = copy.deepcopy(source_metadata)
            copied_metadata["shape_name"] = names_by_new_id[new_id]
            copied_metadata["slide_index"] = index
            if collection == "connectors":
                for endpoint in ("source", "target"):
                    endpoint_id = int(copied_metadata.get(f"{endpoint}_id", 0))
                    if endpoint_id in id_mapping:
                        copied_metadata[f"{endpoint}_id"] = id_mapping[endpoint_id]
                        copied_metadata[f"{endpoint}_name"] = names_by_new_id[id_mapping[endpoint_id]]
            metadata.setdefault(collection, {})[str(new_id)] = copied_metadata
    _save(prs, state, path)
    return {**_shape_result(duplicate, index), "duplicated_from": source.name}


def action_group_shapes(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    _assert_shape_name_available(slide, args.get("name"))
    members = [_shape(slide, {"shape_name": name}) for name in args["shape_names"]]
    group = slide.shapes.add_group_shape(members)
    if args.get("name"):
        group.name = args["name"]
    _save(prs, state, path)
    return {**_shape_result(group, index), "members": [shape.name for shape in group.shapes]}


def action_ungroup_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    group = _shape(slide, args)
    if shape_type_name(group) != "GROUP":
        raise ValueError(f"{group.name} is not a group")
    transform = group._element.grpSpPr.xfrm
    if str(transform.get("flipH", "0")).lower() in {"1", "true"} or str(transform.get("flipV", "0")).lower() in {"1", "true"}:
        raise RuntimeError("Ungrouping a flipped OOXML group is refused because it would corrupt child geometry.")
    if transform.chExt.cx == 0 or transform.chExt.cy == 0:
        raise RuntimeError("Ungrouping is unsafe because the group child coordinate system has zero extent.")
    scale_x = transform.ext.cx / transform.chExt.cx
    scale_y = transform.ext.cy / transform.chExt.cy
    group_center_x = transform.off.x + transform.ext.cx / 2
    group_center_y = transform.off.y + transform.ext.cy / 2
    group_rotation = float(group.rotation or 0)
    angle = math.radians(group_rotation)
    names = []
    for member in list(group.shapes):
        names.append(member.name)
        member_width = member.width * scale_x
        member_height = member.height * scale_y
        member_center_x = transform.off.x + (member.left - transform.chOff.x + member.width / 2) * scale_x
        member_center_y = transform.off.y + (member.top - transform.chOff.y + member.height / 2) * scale_y
        if group_rotation:
            dx, dy = member_center_x - group_center_x, member_center_y - group_center_y
            member_center_x = group_center_x + math.cos(angle) * dx - math.sin(angle) * dy
            member_center_y = group_center_y + math.sin(angle) * dx + math.cos(angle) * dy
        member.left = int(round(member_center_x - member_width / 2))
        member.top = int(round(member_center_y - member_height / 2))
        member.width = int(round(member_width))
        member.height = int(round(member_height))
        member.rotation = float(member.rotation or 0) + group_rotation
        slide.shapes._spTree.insert_element_before(member._element, "p:extLst")
    group._element.getparent().remove(group._element)
    _save(prs, state, path)
    return {"slide_index": index, "ungrouped": args.get("shape_name") or args.get("shape_id"), "members": names}


def shape_type_name(shape) -> str:
    try:
        return shape.shape_type.name
    except Exception:
        return str(shape.shape_type)


def action_set_z_order(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shape = _shape(slide, args)
    parent = shape._element.getparent()
    command = args["command"]
    repeat = int(args.get("repeat", 1))
    for _ in range(repeat):
        siblings = [item for item in parent if item.tag.endswith(("sp", "pic", "graphicFrame", "grpSp", "cxnSp"))]
        pos = siblings.index(shape._element)
        if command == "bring_to_front":
            parent.remove(shape._element)
            parent.insert_element_before(shape._element, "p:extLst")
        elif command == "send_to_back":
            parent.remove(shape._element)
            parent.insert(2, shape._element)
        elif command == "bring_forward" and pos + 1 < len(siblings):
            sibling = siblings[pos + 1]
            parent.remove(shape._element)
            parent.insert(parent.index(sibling) + 1, shape._element)
        elif command == "send_backward" and pos > 0:
            sibling = siblings[pos - 1]
            parent.remove(shape._element)
            parent.insert(parent.index(sibling), shape._element)
    _save(prs, state, path)
    return {**_shape_result(shape, index), "command": command}


def action_align_shapes(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    prs_slide = _slide(prs, index)
    shapes = [_shape(prs_slide, {"shape_name": name}) for name in args["shape_names"]]
    relative_slide = args.get("relative_to") == "slide"
    alignment = args["alignment"]
    if alignment == "left":
        value = 0 if relative_slide else min(s.left for s in shapes)
        for s in shapes: s.left = value
    elif alignment == "right":
        value = prs.slide_width if relative_slide else max(s.left + s.width for s in shapes)
        for s in shapes: s.left = value - s.width
    elif alignment == "center":
        value = prs.slide_width / 2 if relative_slide else sum(s.left + s.width / 2 for s in shapes) / len(shapes)
        for s in shapes: s.left = int(value - s.width / 2)
    elif alignment == "top":
        value = 0 if relative_slide else min(s.top for s in shapes)
        for s in shapes: s.top = value
    elif alignment == "bottom":
        value = prs.slide_height if relative_slide else max(s.top + s.height for s in shapes)
        for s in shapes: s.top = value - s.height
    elif alignment == "middle":
        value = prs.slide_height / 2 if relative_slide else sum(s.top + s.height / 2 for s in shapes) / len(shapes)
        for s in shapes: s.top = int(value - s.height / 2)
    _save(prs, state, path)
    return {"slide_index": index, "alignment": alignment, "shape_names": [s.name for s in shapes]}


def action_distribute_shapes(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    shapes = [_shape(slide, {"shape_name": name}) for name in args["shape_names"]]
    direction = args["direction"]
    if direction == "horizontal":
        shapes.sort(key=lambda s: s.left)
        start = 0 if args.get("relative_to") == "slide" else shapes[0].left
        end = prs.slide_width if args.get("relative_to") == "slide" else shapes[-1].left + shapes[-1].width
        gap = (end - start - sum(s.width for s in shapes)) / (len(shapes) - 1)
        cursor = start
        for shape in shapes:
            shape.left = int(cursor)
            cursor += shape.width + gap
    else:
        shapes.sort(key=lambda s: s.top)
        start = 0 if args.get("relative_to") == "slide" else shapes[0].top
        end = prs.slide_height if args.get("relative_to") == "slide" else shapes[-1].top + shapes[-1].height
        gap = (end - start - sum(s.height for s in shapes)) / (len(shapes) - 1)
        cursor = start
        for shape in shapes:
            shape.top = int(cursor)
            cursor += shape.height + gap
    _save(prs, state, path)
    return {"slide_index": index, "direction": direction, "shape_names": [s.name for s in shapes]}


def action_update_shape(args: dict) -> dict:
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    old_name = shape.name
    if args.get("new_name"):
        _assert_shape_name_available(_slide(prs, index), args["new_name"], shape.shape_id)
        shape.name = args["new_name"]
    for prop in ("left", "top", "width", "height"):
        if args.get(prop) is not None:
            setattr(shape, prop, _pt(args[prop]))
    if args.get("rotation") is not None:
        shape.rotation = float(args["rotation"])
    if args.get("text") is not None:
        _set_text(shape, args["text"], args)
    elif getattr(shape, "has_text_frame", False):
        _apply_text_frame(shape.text_frame, args)
    _set_fill(shape, args.get("fill_color"), args.get("fill_transparency"))
    _set_line(shape, args)
    if args.get("start_arrow") is not None or args.get("end_arrow") is not None:
        if not hasattr(shape, "line"):
            raise ValueError(f"Shape '{shape.name}' does not support arrowhead styling.")
        _set_arrow(
            shape,
            args.get("start_arrow", _get_arrow(shape, "a:headEnd")),
            args.get("end_arrow", _get_arrow(shape, "a:tailEnd")),
        )
    metadata = state.setdefault("metadata", {})
    for collection in ("rasters", "lines", "connectors"):
        entry = metadata.get(collection, {}).get(str(shape.shape_id))
        if entry is not None:
            entry["shape_name"] = shape.name
    if shape.name != old_name:
        for connector in metadata.get("connectors", {}).values():
            if connector.get("source_id") == shape.shape_id:
                connector["source_name"] = shape.name
            if connector.get("target_id") == shape.shape_id:
                connector["target_name"] = shape.name
    line_metadata = metadata.get("lines", {}).get(str(shape.shape_id))
    if line_metadata is not None:
        x1, y1, x2, y2 = _connector_endpoints(shape)
        line_metadata.update({"begin_x": x1, "begin_y": y1, "end_x": x2, "end_y": y2})
        if args.get("start_arrow") is not None:
            line_metadata["start_arrow"] = args["start_arrow"]
        if args.get("end_arrow") is not None:
            line_metadata["end_arrow"] = args["end_arrow"]
    connector_metadata = metadata.get("connectors", {}).get(str(shape.shape_id))
    if connector_metadata is not None:
        if args.get("start_arrow") is not None:
            connector_metadata["start_arrow"] = args["start_arrow"]
        if args.get("end_arrow") is not None:
            connector_metadata["end_arrow"] = args["end_arrow"]
    _save(prs, state, path)
    return _shape_result(shape, index)


def action_delete_shape(args: dict) -> dict:
    if args.get("confirm") is not True:
        raise ValueError("confirm=true is required")
    state, path, prs = _load()
    index = int(args["slide_index"])
    shape = _shape(_slide(prs, index), args)
    deleted = {"shape_name": shape.name, "shape_id": shape.shape_id}
    metadata = state.setdefault("metadata", {})
    for collection in ("rasters", "lines", "connectors"):
        metadata.get(collection, {}).pop(str(shape.shape_id), None)
    shape._element.getparent().remove(shape._element)
    _save(prs, state, path)
    return {"slide_index": index, "deleted": deleted}


def _segment_intersects_rect(x1: float, y1: float, x2: float, y2: float, shape, inset: float) -> bool:
    margin = max(0.0, inset)
    if shape.width.pt <= 2 * margin or shape.height.pt <= 2 * margin:
        margin = 0.0
    left, right = shape.left.pt + margin, shape.left.pt + shape.width.pt - margin
    top, bottom = shape.top.pt + margin, shape.top.pt + shape.height.pt - margin
    if (left < x1 < right and top < y1 < bottom) or (left < x2 < right and top < y2 < bottom):
        return True
    dx, dy = x2 - x1, y2 - y1
    p = (-dx, dx, -dy, dy)
    q = (x1 - left, right - x1, y1 - top, bottom - y1)
    minimum, maximum = 0.0, 1.0
    for direction, distance in zip(p, q):
        if abs(direction) < 1e-9:
            if distance < 0:
                return False
            continue
        ratio = distance / direction
        if direction < 0:
            minimum = max(minimum, ratio)
        else:
            maximum = min(maximum, ratio)
        if minimum > maximum:
            return False
    return maximum - minimum > 1e-6


def _strict_segment_crossing(first: tuple[float, float, float, float], second: tuple[float, float, float, float], clearance: float) -> tuple[float, float] | None:
    ax1, ay1, ax2, ay2 = first
    bx1, by1, bx2, by2 = second
    adx, ady, bdx, bdy = ax2 - ax1, ay2 - ay1, bx2 - bx1, by2 - by1
    denominator = adx * bdy - ady * bdx
    if abs(denominator) < 1e-9:
        return None
    cx, cy = bx1 - ax1, by1 - ay1
    a_ratio = (cx * bdy - cy * bdx) / denominator
    b_ratio = (cx * ady - cy * adx) / denominator
    if not 1e-6 < a_ratio < 1 - 1e-6 or not 1e-6 < b_ratio < 1 - 1e-6:
        return None
    a_length, b_length = math.hypot(adx, ady), math.hypot(bdx, bdy)
    endpoint_distance = min(a_ratio * a_length, (1 - a_ratio) * a_length, b_ratio * b_length, (1 - b_ratio) * b_length)
    if endpoint_distance <= clearance:
        return None
    return ax1 + a_ratio * adx, ay1 + a_ratio * ady


def _is_route_obstacle(shape, slide_area: float) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.LINE or shape.width.pt <= 2 or shape.height.pt <= 2:
        return False
    name_tokens = shape.name.lower().replace("-", "_").split("_")
    if any(token in {"panel", "background", "bg", "container", "region", "frame"} for token in name_tokens):
        return False
    return (shape.width * shape.height) / slide_area < 0.18


def action_audit_figure(args: dict) -> dict:
    state, path, prs = _load(False)
    index = int(args["slide_index"])
    slide = _slide(prs, index)
    findings = []
    max_findings = int(args.get("max_findings", 300))
    alignment_tolerance = float(args.get("alignment_tolerance", 0.75))
    endpoint_clearance = float(args.get("endpoint_clearance", 1.5))
    slide_area = max(1, prs.slide_width * prs.slide_height)
    metadata = state.get("metadata", {})
    raster_meta = metadata.get("rasters", {})
    shapes_by_id = {shape.shape_id: shape for shape in slide.shapes}
    name_counts: dict[str, int] = {}
    for shape in slide.shapes:
        name_counts[shape.name.lower()] = name_counts.get(shape.name.lower(), 0) + 1
        if (shape.width <= 0 or shape.height <= 0) and shape.shape_type != MSO_SHAPE_TYPE.LINE:
            findings.append({"severity": "hard", "category": "geometry", "shape_name": shape.name, "message": "Shape has non-positive dimensions."})
        if shape.left.pt < -alignment_tolerance or shape.top.pt < -alignment_tolerance or shape.left.pt + shape.width.pt > prs.slide_width.pt + alignment_tolerance or shape.top.pt + shape.height.pt > prs.slide_height.pt + alignment_tolerance:
            findings.append({"severity": "hard", "category": "bounds", "shape_name": shape.name, "message": "Shape extends beyond slide bounds."})
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ratio = (shape.width * shape.height) / slide_area
            declaration = raster_meta.get(str(shape.shape_id))
            if declaration is None:
                findings.append({"severity": "hard", "category": "raster_editability", "shape_name": shape.name, "message": "Picture has no atomic-raster declaration."})
            else:
                crop = declaration.get("crop") or {}
                cropped = any(float(value or 0) > 0 for value in crop.values())
                if len(str(declaration.get("raster_reason", "")).strip()) < 8:
                    findings.append({"severity": "hard", "category": "raster_reason", "shape_name": shape.name, "message": "Picture raster reason is missing or too vague."})
                if declaration.get("atomic_raster_unit") is not True or declaration.get("contains_reconstructable_content") is not False:
                    findings.append({"severity": "hard", "category": "raster_atomicity", "shape_name": shape.name, "message": "Picture is not declared as one irreducible raster field."})
                if declaration.get("source_is_tightly_cropped") is not True and not cropped:
                    findings.append({"severity": "hard", "category": "raster_crop", "shape_name": shape.name, "message": "Picture is neither tightly cropped nor cropped through picture properties."})
                if len(str(declaration.get("decomposition_note", "")).strip()) < 8:
                    findings.append({"severity": "hard", "category": "raster_decomposition", "shape_name": shape.name, "message": "Picture has no useful decomposition note."})
            if ratio >= float(args.get("large_raster_area_ratio", 0.08)):
                findings.append({"severity": "warning", "category": "large_raster_surface", "shape_name": shape.name, "message": "Large picture requires full-resolution review for further decomposition."})
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            font_size = 18.0
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        font_size = run.font.size.pt
                        break
            approx_capacity = max(1, (shape.width.pt / max(font_size * 0.55, 1)) * (shape.height.pt / max(font_size * 1.25, 1)))
            if len(shape.text) > approx_capacity * 1.3:
                findings.append({"severity": "warning", "category": "text_fit", "shape_name": shape.name, "message": "Text may overflow; verify the rendered slide."})
        if len(findings) >= max_findings:
            break

    for name, count in name_counts.items():
        if count > 1 and len(findings) < max_findings:
            findings.append({"severity": "hard", "category": "duplicate_name", "shape_name": name, "message": f"Semantic shape name occurs {count} times."})

    obstacles = [shape for shape in slide.shapes if _is_route_obstacle(shape, slide_area)]
    straight_segments = []
    for shape_id, line_info in metadata.get("lines", {}).items():
        if int(line_info.get("slide_index", 0)) != index or int(shape_id) not in shapes_by_id:
            continue
        shape = shapes_by_id[int(shape_id)]
        segment = _connector_endpoints(shape)
        straight_segments.append({"shape": shape, "segment": segment, "endpoints": set()})
        for obstacle in obstacles:
            if obstacle.shape_id == shape.shape_id:
                continue
            if _segment_intersects_rect(*segment, obstacle, endpoint_clearance) and len(findings) < max_findings:
                findings.append({"severity": "hard", "category": "connector_path_through_object", "shape_name": shape.name, "message": f"Straight line passes through unrelated object '{obstacle.name}'."})
        for label, x, y, arrow in (
            ("start", segment[0], segment[1], line_info.get("start_arrow", "none")),
            ("end", segment[2], segment[3], line_info.get("end_arrow", "none")),
        ):
            if arrow == "none":
                continue
            for obstacle in obstacles:
                if obstacle.shape_id == shape.shape_id:
                    continue
                if obstacle.left.pt + endpoint_clearance < x < obstacle.left.pt + obstacle.width.pt - endpoint_clearance and obstacle.top.pt + endpoint_clearance < y < obstacle.top.pt + obstacle.height.pt - endpoint_clearance:
                    if len(findings) < max_findings:
                        findings.append({"severity": "hard", "category": "arrowhead_intrusion", "shape_name": shape.name, "message": f"{label} arrowhead lies inside '{obstacle.name}'."})
                    break

    for shape_id, connector_info in metadata.get("connectors", {}).items():
        if int(connector_info.get("slide_index", 0)) != index or int(shape_id) not in shapes_by_id:
            continue
        connector = shapes_by_id[int(shape_id)]
        source = shapes_by_id.get(int(connector_info.get("source_id", 0)))
        target = shapes_by_id.get(int(connector_info.get("target_id", 0)))
        properties = connector._element.nvCxnSpPr.cNvCxnSpPr
        start_binding = properties.find("a:stCxn", properties.nsmap)
        end_binding = properties.find("a:endCxn", properties.nsmap)
        binding_valid = (
            source is not None and target is not None and start_binding is not None and end_binding is not None
            and start_binding.get("id") == str(source.shape_id) and end_binding.get("id") == str(target.shape_id)
            and start_binding.get("idx") == str(int(connector_info["source_site"]) - 1)
            and end_binding.get("idx") == str(int(connector_info["target_site"]) - 1)
        )
        if not binding_valid and len(findings) < max_findings:
            findings.append({"severity": "hard", "category": "connector_detached", "shape_name": connector.name, "message": "Connector source/target connection-site binding is missing or stale."})
        if source is None or target is None or connector_info.get("connector_type") != "straight":
            continue
        segment = (*_site_point(source, int(connector_info["source_site"])), *_site_point(target, int(connector_info["target_site"])))
        endpoint_names = {source.name.lower(), target.name.lower()}
        straight_segments.append({"shape": connector, "segment": segment, "endpoints": endpoint_names})
        for obstacle in obstacles:
            if obstacle.shape_id in {connector.shape_id, source.shape_id, target.shape_id}:
                continue
            if _segment_intersects_rect(*segment, obstacle, endpoint_clearance) and len(findings) < max_findings:
                findings.append({"severity": "hard", "category": "connector_path_through_object", "shape_name": connector.name, "message": f"Straight connector passes through unrelated object '{obstacle.name}'."})

    for first_index, first in enumerate(straight_segments):
        for second in straight_segments[first_index + 1:]:
            if first["endpoints"].intersection(second["endpoints"]):
                continue
            crossing = _strict_segment_crossing(first["segment"], second["segment"], endpoint_clearance)
            if crossing is not None and len(findings) < max_findings:
                findings.append({
                    "severity": "hard",
                    "category": "connector_crossing",
                    "shape_name": first["shape"].name,
                    "message": f"Route crosses '{second['shape'].name}' near ({crossing[0]:.2f}, {crossing[1]:.2f}).",
                })
    hard_count = sum(1 for item in findings if item["severity"] == "hard")
    return {
        "slide_index": index,
        "path": str(path),
        "backend": "python-pptx-ooxml+application-reload",
        "host_application": state.get("host_application", "auto"),
        "shape_count": len(slide.shapes),
        "findings": findings,
        "hard_failure_count": hard_count,
        "warning_count": len(findings) - hard_count,
        "passed_deterministic_gate": hard_count == 0,
        "renderer_review_required": True,
    }


def _find_binary(name: str) -> str:
    binary = shutil.which(name)
    if not binary:
        raise RuntimeError(f"Required local renderer binary not found: {name}")
    return binary


def _render_pdf(source: Path, destination: Path) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="scientific-illustrator-pdf-") as tmp:
        tmp_path = Path(tmp)
        _run([_find_binary("soffice"), "--headless", "--convert-to", "pdf", "--outdir", str(tmp_path), str(source)], timeout=180)
        generated = tmp_path / f"{source.stem}.pdf"
        if not generated.exists():
            candidates = list(tmp_path.glob("*.pdf"))
            if not candidates:
                raise RuntimeError("LibreOffice did not produce a PDF.")
            generated = candidates[0]
        shutil.copy2(generated, destination)


def _render_dimensions(path: Path, args: dict) -> tuple[int, int, int, int, bool]:
    presentation = Presentation(path)
    _slide(presentation, int(args["slide_index"]))
    requested_width = int(args.get("width", 1920))
    height_was_explicit = args.get("height") is not None
    requested_height = int(args["height"]) if height_was_explicit else max(1, round(requested_width * presentation.slide_height / presentation.slide_width))
    preserve = args.get("preserve_aspect_ratio", True) is not False
    if not preserve:
        return requested_width, requested_height, requested_width, requested_height, False
    if not height_was_explicit:
        return requested_width, requested_height, requested_width, requested_height, True
    scale = min(requested_width / presentation.slide_width, requested_height / presentation.slide_height)
    width = max(1, round(presentation.slide_width * scale))
    height = max(1, round(presentation.slide_height * scale))
    return requested_width, requested_height, width, height, True


def action_export_slide_image(args: dict) -> dict:
    state, path = _require_path(False)
    output = Path(args["output_path"]).expanduser().resolve()
    if output.exists() and not args.get("overwrite", False):
        raise FileExistsError(f"Output exists: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    requested_width, requested_height, width, height, preserved = _render_dimensions(path, args)
    slide_index = int(args["slide_index"])
    with tempfile.TemporaryDirectory(prefix="scientific-illustrator-render-") as tmp:
        pdf_path = Path(tmp) / "deck.pdf"
        _render_pdf(path, pdf_path)
        jpeg = output.suffix.lower() in (".jpg", ".jpeg")
        rendered = Path(tmp) / ("slide.jpg" if jpeg else "slide.png")
        format_flag = "-jpeg" if jpeg else "-png"
        _run([_find_binary("pdftoppm"), "-f", str(slide_index), "-l", str(slide_index), "-singlefile", format_flag, "-scale-to-x", str(width), "-scale-to-y", str(height), str(pdf_path), str(rendered.with_suffix(""))], timeout=180)
        shutil.copy2(rendered, output)
        mime = "image/jpeg" if jpeg else "image/png"
    return {
        "slide_index": slide_index,
        "output_path": str(output),
        "requested_width": requested_width,
        "requested_height": requested_height,
        "width": width,
        "height": height,
        "aspect_ratio_preserved": preserved,
        "mime_type": mime,
        "renderer": "LibreOffice+Poppler",
    }


def action_save(args: dict) -> dict:
    state, path = _require_path(False)
    output_raw = args.get("output_path")
    if not output_raw:
        if args.get("format") not in (None, "pptx"):
            raise ValueError("PDF export requires an absolute .pdf output_path.")
        return {"output_path": str(path), "format": "pptx", "saved": True}
    output = Path(output_raw).expanduser().resolve()
    if output.exists() and not args.get("overwrite", False):
        raise FileExistsError(f"Output exists: {output}")
    fmt = args.get("format") or ("pdf" if output.suffix.lower() == ".pdf" else "pptx")
    expected_suffix = ".pdf" if fmt == "pdf" else ".pptx"
    if output.suffix.lower() != expected_suffix:
        raise ValueError(f"format '{fmt}' requires an {expected_suffix} output_path.")
    output.parent.mkdir(parents=True, exist_ok=True)
    if fmt == "pdf":
        _render_pdf(path, output)
    else:
        shutil.copy2(path, output)
        state["last_saved_output"] = str(output)
        _write_state(state)
    return {"output_path": str(output), "format": fmt, "saved": True}


def action_close_presentation(args: dict) -> dict:
    if args.get("confirm") is not True:
        raise ValueError("confirm=true is required")
    state, path = _require_path(False)
    if args.get("save_changes") == "save" and args.get("output_path"):
        action_save({"output_path": args["output_path"], "overwrite": args.get("overwrite", False), "format": "pptx"})
    host_name, _ = _select_host(state=state)
    if sys.platform != "darwin" or host_name != "powerpoint":
        return {
            "closed": False,
            "detached": True,
            "path": str(path),
            "file_retained": path.exists(),
            "note": "The file-backed WPS backend does not force-close application windows; close the managed deck in WPS when convenient.",
        }
    available = _available_hosts()
    powerpoint_host = available["powerpoint"]
    process_ids = _main_process_ids("powerpoint", powerpoint_host)
    document_open_before = _document_open_verification(path, process_ids, "powerpoint")
    if document_open_before is False:
        return {
            "closed": True,
            "already_closed": True,
            "close_verified": True,
            "document_open_before": False,
            "document_open_verified": False,
            "path": str(path),
            "file_retained": path.exists(),
        }
    script = r'''
on run argv
    set targetDirectory to item 1 of argv
    set resolvedDirectory to item 2 of argv
    set targetName to item 3 of argv
    tell application "Microsoft PowerPoint"
        repeat with deckIndex from 1 to count of presentations
            try
                if (get name of presentation deckIndex) is (my targetName) then
                    if ((get path of presentation deckIndex) as text) is (my targetDirectory) or ((get path of presentation deckIndex) as text) is (my resolvedDirectory) then close presentation deckIndex saving no
                end if
            end try
        end repeat
    end tell
end run
'''
    close_timed_out = False
    close_return_code = None
    close_error = None
    try:
        close_attempt = _run([
            "osascript", "-e", script,
            str(path.parent), str(path.parent.resolve()), path.name,
        ], check=False, timeout=8)
        close_return_code = close_attempt.returncode
        close_error = close_attempt.stderr.strip() or None
    except subprocess.TimeoutExpired:
        close_timed_out = True
        close_error = "PowerPoint did not answer the exact-file close request within 8 seconds."
    process_ids = _main_process_ids("powerpoint", powerpoint_host)
    document_open = _document_open_verification(path, process_ids, "powerpoint")
    closed = document_open is False
    return {
        "closed": closed,
        "close_verified": closed,
        "close_timed_out": close_timed_out,
        "close_return_code": close_return_code,
        "close_error": close_error,
        "document_open_before": document_open_before,
        "document_open_verified": document_open,
        "path": str(path),
        "file_retained": path.exists(),
    }


def action_quit_application(args: dict) -> dict:
    if args.get("confirm") is not True:
        raise ValueError("confirm=true is required")
    state = _state()
    host = _presentation_host_info(args, state)
    if host["host_application"] != "powerpoint" or sys.platform != "darwin":
        raise RuntimeError("Safe programmatic application quit is only available for Microsoft PowerPoint on macOS in the file-backed backend.")
    if host["presentation_count"] is None:
        raise RuntimeError("PowerPoint presentation count could not be verified; refusing to quit the application.")
    if host["presentation_count"] != 0:
        raise RuntimeError("PowerPoint still has open presentations; close them explicitly before quitting.")
    expected = int(args.get("expected_process_id", 0))
    if expected not in host["process_ids"]:
        raise RuntimeError("expected_process_id does not match the current PowerPoint process.")
    _osascript('tell application "Microsoft PowerPoint" to quit', check=True)
    return {"quit": True, "process_id": expected}


ACTIONS = {
    "status": action_status,
    "capabilities": action_capabilities,
    "launch": action_launch,
    "new_presentation": action_new_presentation,
    "inspect": action_inspect,
    "audit_figure": action_audit_figure,
    "activate_slide": action_activate_slide,
    "refresh": action_refresh,
    "add_slide": action_add_slide,
    "add_textbox": action_add_textbox,
    "add_shape": action_add_shape,
    "add_image": action_add_image,
    "add_line": action_add_line,
    "add_connector": action_add_connector,
    "add_table": action_add_table,
    "update_table_cell": action_update_table_cell,
    "update_table_layout": action_update_table_layout,
    "add_chart": action_add_chart,
    "duplicate_shape": action_duplicate_shape,
    "group_shapes": action_group_shapes,
    "ungroup_shape": action_ungroup_shape,
    "set_z_order": action_set_z_order,
    "align_shapes": action_align_shapes,
    "distribute_shapes": action_distribute_shapes,
    "update_shape": action_update_shape,
    "delete_shape": action_delete_shape,
    "export_slide_image": action_export_slide_image,
    "save": action_save,
    "close_presentation": action_close_presentation,
    "quit_application": action_quit_application,
}


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("Usage: powerpoint-mac-bridge.py <payload-base64>")
    payload = json.loads(base64.b64decode(sys.argv[1]).decode("utf-8"))
    action = str(payload["action"])
    arguments = payload.get("arguments") or {}
    if action not in ACTIONS:
        raise ValueError(f"Unsupported macOS PowerPoint action: {action}")
    result = ACTIONS[action](arguments)
    sys.stdout.write(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        sys.stderr.write(f"{type(exc).__name__}: {exc}\n")
        raise
